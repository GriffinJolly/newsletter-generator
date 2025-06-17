import os
import json
import re
from collections import defaultdict, Counter
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from datetime import datetime
import calendar

def clean_filename(name):
    return re.sub(r'\W+', '_', name).strip("_")

def extract_company_name(article):
    """Extract company name from article - improved logic"""
    title = article.get("title", "").lower()
    
    # Check for Slaughter and May variations
    if "slaughter" in title and "may" in title:
        return "Slaughter and May"
    
    # Check in summary as backup
    summary = article.get("summary", "").lower()
    if "slaughter" in summary and "may" in summary:
        return "Slaughter and May"
    
    # If not found, check if it's a competitor article about law firms
    if article.get("company_type") == "competitor":
        return "Slaughter and May"  # Assuming all competitor articles are about S&M
    
    return "Unknown Company"

def analyze_content_themes(articles):
    """Analyze articles to extract key themes and topics"""
    themes = defaultdict(int)
    sentiment_keywords = {
        'positive': ['success', 'growth', 'expansion', 'award', 'wins', 'appointed', 'promotion', 'innovative', 'leading', 'excellence'],
        'negative': ['challenge', 'decline', 'issue', 'problem', 'concern', 'risk', 'loss', 'departure', 'difficult'],
        'neutral': ['announcement', 'change', 'update', 'move', 'development', 'report']
    }
    
    sentiment_scores = {'positive': 0, 'negative': 0, 'neutral': 0}
    
    for article in articles:
        text = (article.get('title', '') + ' ' + article.get('summary', '')).lower()
        
        # Count sentiment indicators
        for sentiment, keywords in sentiment_keywords.items():
            for keyword in keywords:
                if keyword in text:
                    sentiment_scores[sentiment] += 1
        
        # Extract potential themes (simple keyword extraction)
        common_legal_terms = ['merger', 'acquisition', 'deal', 'transaction', 'advisory', 'counsel', 'litigation', 
                             'compliance', 'regulation', 'partnership', 'client', 'finance', 'corporate', 'banking']
        
        for term in common_legal_terms:
            if term in text:
                themes[term] += 1
    
    return themes, sentiment_scores

def calculate_media_metrics(articles):
    """Calculate various media coverage metrics"""
    sources = [article.get('source', 'Unknown') for article in articles]
    source_diversity = len(set(sources))
    
    # Calculate publication frequency
    dates = [article.get('published', '') for article in articles if article.get('published', '')]
    date_spread = 0
    if len(dates) > 1:
        try:
            sorted_dates = sorted(dates)
            # Simple date spread calculation (could be enhanced)
            date_spread = len(set(date[:7] for date in sorted_dates))  # Count unique months
        except:
            pass
    
    categories = [article.get('category', 'News') for article in articles]
    category_distribution = Counter(categories)
    
    return {
        'source_diversity': source_diversity,
        'date_spread_months': date_spread,
        'category_distribution': category_distribution,
        'total_sources': sources
    }

def generate_strategic_insights(company_name, articles, themes, sentiment_scores, metrics):
    """Generate strategic insights based on analysis"""
    insights = []
    recommendations = []
    
    # Media presence analysis
    total_articles = len(articles)
    if total_articles > 50:
        insights.append(f"🎯 Strong Media Presence: {company_name} maintains high visibility with {total_articles} articles, indicating active market engagement")
    elif total_articles > 20:
        insights.append(f"📊 Moderate Media Coverage: {total_articles} articles suggest steady but focused media attention")
    else:
        insights.append(f"📈 Emerging Coverage: {total_articles} articles indicate opportunity for increased media presence")
    
    # Source diversity analysis
    source_diversity = metrics['source_diversity']
    if source_diversity > 10:
        insights.append(f"🌐 Excellent Source Diversification: Coverage spans {source_diversity} different publications, ensuring broad reach")
    elif source_diversity > 5:
        insights.append(f"📰 Good Media Reach: Featured across {source_diversity} publications with room for expansion")
    else:
        insights.append(f"🎯 Concentrated Coverage: Limited to {source_diversity} sources - opportunity for broader media engagement")
    
    # Sentiment analysis
    total_sentiment = sum(sentiment_scores.values())
    if total_sentiment > 0:
        positive_ratio = sentiment_scores['positive'] / total_sentiment
        if positive_ratio > 0.6:
            insights.append("✅ Predominantly Positive Sentiment: Media coverage reflects favorably on company activities and reputation")
        elif positive_ratio > 0.4:
            insights.append("⚖️ Balanced Media Tone: Mix of positive and neutral coverage suggests authentic reporting")
        else:
            insights.append("⚠️ Mixed Sentiment Indicators: Opportunity to enhance positive narrative in media coverage")
    
    # Theme analysis
    top_themes = sorted(themes.items(), key=lambda x: x[1], reverse=True)[:3]
    if top_themes:
        theme_text = ", ".join([f"{theme} ({count})" for theme, count in top_themes])
        insights.append(f"🔍 Key Coverage Themes: {theme_text} - indicating primary areas of market focus")
    
    # Temporal analysis
    if metrics['date_spread_months'] > 6:
        insights.append(f"📅 Sustained Coverage: Articles span {metrics['date_spread_months']} months, showing consistent media attention")
    elif metrics['date_spread_months'] > 3:
        insights.append(f"⏱️ Regular Coverage: {metrics['date_spread_months']} months of coverage indicates ongoing market relevance")
    
    # Generate recommendations
    recommendations.append("📈 Strategic Recommendations:")
    
    if source_diversity < 8:
        recommendations.append("• Expand media outreach to additional industry publications and mainstream business media")
    
    if sentiment_scores['positive'] < sentiment_scores['negative']:
        recommendations.append("• Develop proactive media strategy to enhance positive narrative and thought leadership")
    
    if 'innovation' not in themes or themes.get('innovation', 0) < 3:
        recommendations.append("• Increase coverage of innovative practices and technological advancements")
    
    recommendations.append("• Leverage high-performing content themes for future PR initiatives")
    recommendations.append("• Monitor competitor coverage patterns to identify market positioning opportunities")
    recommendations.append("• Establish regular media cadence for consistent visibility and relationship building")
    
    # Market positioning insights
    category_dist = metrics['category_distribution']
    if len(category_dist) > 3:
        insights.append(f"🎪 Diverse Content Portfolio: Coverage across {len(category_dist)} categories demonstrates market versatility")
    
    return insights, recommendations

def add_gradient_background(slide, color1=(240, 248, 255), color2=(176, 196, 222)):
    """Add a subtle gradient background to slide"""
    try:
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*color1)
    except Exception as e:
        # Fallback to simple solid background if gradient fails
        print(f"Warning: Could not apply gradient background, using solid color instead")

def create_summary_slide(prs, company_name, articles):
    """Create an executive summary slide with visualizations"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    add_gradient_background(slide, (250, 250, 255), (230, 230, 250))
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "Executive Summary"
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(25, 25, 112)  # Midnight blue
    title_para.alignment = PP_ALIGN.CENTER
    
    # Calculate statistics safely
    published_dates = [article.get('published', '') for article in articles if article.get('published', '')]
    date_range = "N/A"
    if published_dates:
        try:
            min_date = min(published_dates)
            max_date = max(published_dates)
            date_range = f"{min_date} to {max_date}"
        except:
            date_range = "Various dates"
    
    # Statistics boxes
    stats = [
        ("📰 Total Articles", len(articles)),
        ("📅 Date Range", date_range),
        ("Categories", len(set(article.get('category', 'News') for article in articles))),
        ("Sources", len(set(article.get('source', 'Unknown') for article in articles)))
    ]
    
    # Create stat boxes in a 2x2 grid
    box_width, box_height = Inches(4), Inches(1.2)
    positions = [(Inches(0.5), Inches(2)), (Inches(5.5), Inches(2)),
                 (Inches(0.5), Inches(3.5)), (Inches(5.5), Inches(3.5))]
    
    for i, (stat_name, stat_value) in enumerate(stats):
        if i < len(positions):
            left, top = positions[i]
            
            # Create rounded rectangle shape
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, box_width, box_height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(240, 248, 255)  # Alice blue
            shape.line.color.rgb = RGBColor(70, 130, 180)  # Steel blue
            shape.line.width = Pt(2)
            
            # Add text
            text_frame = shape.text_frame
            text_frame.margin_left = Inches(0.1)
            text_frame.margin_right = Inches(0.1)
            text_frame.margin_top = Inches(0.1)
            text_frame.margin_bottom = Inches(0.1)
            
            p1 = text_frame.paragraphs[0]
            p1.text = stat_name
            p1.font.size = Pt(12)
            p1.font.bold = True
            p1.font.color.rgb = RGBColor(25, 25, 112)
            p1.alignment = PP_ALIGN.CENTER
            
            p2 = text_frame.add_paragraph()
            p2.text = str(stat_value)
            p2.font.size = Pt(14)
            p2.font.bold = True
            p2.font.color.rgb = RGBColor(0, 100, 0)  # Dark green
            p2.alignment = PP_ALIGN.CENTER

def create_chart_slide(prs, company_name, articles):
    """Create a slide with category distribution chart"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    add_gradient_background(slide, (255, 250, 250), (255, 228, 225))
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(0.8))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "📈 Article Distribution by Category"
    title_para.font.size = Pt(24)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(139, 0, 0)  # Dark red
    title_para.alignment = PP_ALIGN.CENTER
    
    # Count articles by category
    category_counts = Counter(article.get('category', 'News') for article in articles)
    
    # Only create chart if we have data
    if category_counts:
        try:
            # Create chart data
            chart_data = CategoryChartData()
            chart_data.categories = list(category_counts.keys())
            chart_data.add_series('Articles', list(category_counts.values()))
            
            # Add chart
            x, y, cx, cy = Inches(1.5), Inches(1.5), Inches(7), Inches(4.5)
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
            ).chart
            
            # Customize chart
            chart.has_legend = False
            #chart.legend.position = XL_LEGEND_POSITION.RIGHT
            #chart.plots[0].has_data_labels = True
        except Exception as e:
            # Fallback: Create a simple text summary instead of chart
            summary_box = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(7), Inches(4))
            summary_frame = summary_box.text_frame
            summary_frame.text = "Category Distribution:\n\n"
            for category, count in category_counts.items():
                summary_frame.text += f"• {category}: {count} articles\n"

def create_timeline_slide(prs, company_name, articles):
    """Create a timeline visualization slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    add_gradient_background(slide, (248, 255, 248), (240, 255, 240))
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(8), Inches(0.8))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "📅 Publication Timeline"
    title_para.font.size = Pt(24)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(0, 100, 0)  # Dark green
    title_para.alignment = PP_ALIGN.CENTER
    
    # Process dates and create monthly distribution
    monthly_counts = defaultdict(int)
    for article in articles:
        date_str = article.get('published', '')
        if date_str:
            try:
                # Assuming date format is YYYY-MM-DD or similar
                month_key = date_str[:7]  # YYYY-MM
                monthly_counts[month_key] += 1
            except:
                pass
    
    if monthly_counts:
        try:
            # Create bar chart data
            chart_data = CategoryChartData()
            sorted_months = sorted(monthly_counts.keys())
            chart_data.categories = [f"{month.split('-')[1]}/{month.split('-')[0]}" for month in sorted_months]
            chart_data.add_series('Articles Published', [monthly_counts[month] for month in sorted_months])
            
            # Add chart
            x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(4.5)
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
            ).chart
            
            chart.has_legend = False
            chart.plots[0].has_data_labels = True
        except Exception as e:
            # Fallback: Create a simple text timeline instead of chart
            timeline_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
            timeline_frame = timeline_box.text_frame
            timeline_frame.text = "Publication Timeline:\n\n"
            sorted_months = sorted(monthly_counts.keys())
            for month in sorted_months:
                timeline_frame.text += f"• {month}: {monthly_counts[month]} articles\n"
    else:
        # No date data available
        no_data_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(2))
        no_data_frame = no_data_box.text_frame
        no_data_frame.text = "📅 No publication date information available for timeline analysis"
        no_data_frame.paragraphs[0].font.size = Pt(16)
        no_data_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def create_enhanced_conclusion_slide(prs, company_name, articles):
    """Create an enhanced conclusion slide with detailed insights and recommendations"""
    conclusion_slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_gradient_background(conclusion_slide, (255, 250, 240), (255, 245, 238))
    
    # Analyze the articles for insights
    themes, sentiment_scores = analyze_content_themes(articles)
    metrics = calculate_media_metrics(articles)
    insights, recommendations = generate_strategic_insights(company_name, articles, themes, sentiment_scores, metrics)
    
    # Define color scheme
    primary_color = RGBColor(25, 25, 112)  # Midnight blue
    secondary_color = RGBColor(70, 130, 180)  # Steel blue
    accent_color = RGBColor(220, 20, 60)  # Crimson
    
    # Title
    conclusion_title = conclusion_slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
    conclusion_title.text_frame.text = f"🎯 Strategic Media Analysis: {company_name}"
    conclusion_title.text_frame.paragraphs[0].font.size = Pt(24)
    conclusion_title.text_frame.paragraphs[0].font.bold = True
    conclusion_title.text_frame.paragraphs[0].font.color.rgb = primary_color
    conclusion_title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Create two columns: Insights and Recommendations
    
    # Key Insights (top)
    insights_box = conclusion_slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(8.6), Inches(2.1))
    insights_frame = insights_box.text_frame
    insights_frame.margin_left = Inches(0.2)
    insights_frame.margin_right = Inches(0.1)
    insights_frame.margin_top = Inches(0.1)

    # Insights header
    insights_header = insights_frame.paragraphs[0]
    insights_header.text = "📊 KEY INSIGHTS"
    insights_header.font.size = Pt(16)
    insights_header.font.bold = True
    insights_header.font.color.rgb = accent_color
    insights_header.space_after = Pt(12)

    # Add insights with better formatting
    for insight in insights:
        p = insights_frame.add_paragraph()
        p.text = insight
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(40, 40, 40)
        p.space_after = Pt(8)
        p.line_spacing = 1.2

    # Strategic Recommendations (below insights)
    recommendations_box = conclusion_slide.shapes.add_textbox(Inches(0.7), Inches(3.45), Inches(8.6), Inches(2.1))
    recommendations_frame = recommendations_box.text_frame
    recommendations_frame.margin_left = Inches(0.1)
    recommendations_frame.margin_right = Inches(0.2)
    recommendations_frame.margin_top = Inches(0.1)

    # Recommendations header
    rec_header = recommendations_frame.paragraphs[0]
    rec_header.text = f"{recommendations[0]}"  # "📈 Strategic Recommendations:"
    rec_header.font.size = Pt(16)
    rec_header.font.bold = True
    rec_header.font.color.rgb = accent_color
    rec_header.space_after = Pt(12)

    # Add recommendations with better formatting
    for rec in recommendations[1:]:  # Skip the header
        p = recommendations_frame.add_paragraph()
        p.text = rec
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(40, 40, 40)
        p.space_after = Pt(8)
        p.line_spacing = 1.2
    
    # Add performance metrics summary at the bottom
    metrics_box = conclusion_slide.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(9), Inches(0.8))
    metrics_frame = metrics_box.text_frame
    
    # Create metrics summary
    total_sentiment = sum(sentiment_scores.values())
    sentiment_ratio = f"{sentiment_scores['positive']}:{sentiment_scores['negative']}:{sentiment_scores['neutral']}" if total_sentiment > 0 else "N/A"
    top_theme = max(themes.items(), key=lambda x: x[1])[0] if themes else "General Coverage"
    
    metrics_text = f"📈 Coverage Metrics: {len(articles)} articles • {metrics['source_diversity']} sources • " \
                  f"{len(metrics['category_distribution'])} categories • Sentiment Ratio (P:N:Nu) {sentiment_ratio} • " \
                  f"Primary Theme: {top_theme.title()}"
    
    metrics_para = metrics_frame.paragraphs[0]
    metrics_para.text = metrics_text
    metrics_para.font.size = Pt(10)
    metrics_para.font.color.rgb = secondary_color
    metrics_para.alignment = PP_ALIGN.CENTER
    
    # Add decorative border
    border_shape = conclusion_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.2), Inches(1), Inches(9.6), Inches(5.0)
    )
    border_shape.fill.background()
    border_shape.line.color.rgb = secondary_color
    border_shape.line.width = Pt(2)
    
    # Send shape to back
    border_shape._element.getparent().remove(border_shape._element)
    conclusion_slide.shapes._spTree.insert(2, border_shape._element)

def create_ppt_for_company(company_name, articles, output_dir):
    prs = Presentation()
    
    # Define color scheme
    primary_color = RGBColor(25, 25, 112)  # Midnight blue
    secondary_color = RGBColor(70, 130, 180)  # Steel blue
    accent_color = RGBColor(220, 20, 60)  # Crimson
    
    # Create title slide with enhanced design
    title_slide_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_slide_layout)
    add_gradient_background(title_slide, (240, 248, 255), (230, 230, 250))
    
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    
    title.text = f"{company_name} News Report"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    subtitle.text = f"Generated on {datetime.today().strftime('%B %d, %Y')}\n📰 {len(articles)} Articles Analyzed\n🔍 Comprehensive Media Coverage Analysis"
    subtitle.text_frame.paragraphs[0].font.size = Pt(18)
    subtitle.text_frame.paragraphs[0].font.color.rgb = secondary_color
    
    # Add company logo placeholder (decorative shape)
    logo_shape = title_slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(8.5), Inches(1), Inches(1), Inches(1)
    )
    logo_shape.fill.solid()
    logo_shape.fill.fore_color.rgb = accent_color
    logo_shape.line.color.rgb = primary_color
    logo_shape.line.width = Pt(3)
    
    # Add summary slides
    create_summary_slide(prs, company_name, articles)
    create_chart_slide(prs, company_name, articles)
    create_timeline_slide(prs, company_name, articles)
    
    # Define bullet slide layout
    bullet_slide_layout = prs.slide_layouts[1]

    # Sort articles by date (newest first)
    sorted_articles = sorted(articles, key=lambda x: x.get('published', ''), reverse=True)

    for i, article in enumerate(sorted_articles, 1):
        slide = prs.slides.add_slide(bullet_slide_layout)
        add_gradient_background(slide, (255, 255, 250), (250, 250, 240))
        shapes = slide.shapes

        # Enhanced title with progress indicator
        slide.shapes.title.text = f"📰 {article.get('category', 'News')} • Article {i} of {len(articles)}"
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(20)
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = primary_color

        # Content with improved styling
        content = shapes.placeholders[1]
        text_frame = content.text_frame
        text_frame.clear()
        text_frame.margin_left = Inches(0.3)
        text_frame.margin_right = Inches(0.3)

        # Article title with better formatting
        p = text_frame.paragraphs[0]
        p.text = f"📌 {article['title']}"
        p.font.bold = True
        p.font.size = Pt(16)
        p.font.color.rgb = primary_color
        p.space_after = Pt(12)

        # Summary with enhanced styling
        summary_para = text_frame.add_paragraph()
        summary_para.text = f"📝 Summary:\n{article['summary']}"
        summary_para.font.size = Pt(13)
        summary_para.font.color.rgb = RGBColor(60, 60, 60)
        summary_para.level = 0
        summary_para.space_after = Pt(10)

        # Metadata section with icons and better organization
        metadata_para = text_frame.add_paragraph()
        metadata_text = f"📅 Published: {article['published']}\n📰 Source: {article['source']}"
        metadata_para.text = metadata_text
        metadata_para.font.size = Pt(11)
        metadata_para.font.color.rgb = RGBColor(100, 100, 100)
        metadata_para.level = 0
        metadata_para.space_after = Pt(8)

        # Link with call-to-action styling
        link_para = text_frame.add_paragraph()
        link_para.text = f"🔗 Read Full Article: {article['link']}"
        link_para.font.size = Pt(10)
        link_para.font.color.rgb = RGBColor(0, 100, 200)
        link_para.font.italic = True
        link_para.level = 0
        
        # Add decorative elements
        if i % 1 == 0:  # Add accent shape every slide
            accent_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(6.5), Inches(1), Inches(0.3)
            )
            accent_shape.fill.solid()
            accent_shape.fill.fore_color.rgb = accent_color
            accent_shape.line.color.rgb = primary_color
            accent_shape.text_frame.text = f"{i}/{len(articles)}"
            accent_shape.text_frame.paragraphs[0].font.size = Pt(10)
            accent_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            accent_shape.text_frame.paragraphs[0].font.bold = True
            accent_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Add enhanced conclusion slide
    create_enhanced_conclusion_slide(prs, company_name, articles)

    # Save the presentation
    output_path = os.path.join(output_dir, f"{clean_filename(company_name)}_consolidated_report.pptx")
    
    try:
        prs.save(output_path)
        print(f"Successfully saved PowerPoint report for {company_name}")
        print(f"Location: {output_path}")
        print(f"Total articles: {len(articles)}")
        print(f"Total slides: {len(prs.slides)} (including title slide)")
    except PermissionError:
        print(f"Error: Could not save {output_path}")
        print("Please close any open PowerPoint files and try again")
    except Exception as e:
        print(f"Unexpected error: {str(e)}")

def generate_reports(json_file_path, output_dir="outputs/ppt"):
    """Generate consolidated report for Slaughter and May"""
    os.makedirs(output_dir, exist_ok=True)

    try:
        with open(json_file_path, "r", encoding="utf-8") as f:
            articles = json.load(f)
        print(f"Loaded {len(articles)} articles from {json_file_path}")
    except FileNotFoundError:
        print(f"Error: Could not find file {json_file_path}")
        return
    except json.JSONDecodeError:
        print(f"Error: Invalid JSON in {json_file_path}")
        return

    # Group articles by company
    grouped_articles = defaultdict(list)
    for article in articles:
        company = extract_company_name(article)
        grouped_articles[company].append(article)

    print(f"Found companies: {list(grouped_articles.keys())}")
    
    # Create PPT for each company
    for company, company_articles in grouped_articles.items():
        print(f"\nProcessing {company}...")
        create_ppt_for_company(company, company_articles, output_dir)

    print(f"\nReport generation complete!")

# Alternative function if you want to force all articles into one company
def generate_single_company_report(json_file_path, company_name="Slaughter and May", output_dir="output_ppts"):
    """Generate a single consolidated report for specified company"""
    os.makedirs(output_dir, exist_ok=True)

    try:
        with open(json_file_path, "r", encoding="utf-8") as f:
            articles = json.load(f)
        print(f"Loaded {len(articles)} articles from {json_file_path}")
    except FileNotFoundError:
        print(f"Error: Could not find file {json_file_path}")
        return
    except json.JSONDecodeError:
        print(f"Error: Invalid JSON in {json_file_path}")
        return

    print(f"Creating consolidated report for {company_name}...")
    create_ppt_for_company(company_name, articles, output_dir)
    print(f"Single company report generation complete!")

# Run this script
if __name__ == "__main__":
    print("Choose an option:")
    print("1. Auto-detect companies from articles (recommended)")
    print("2. Force all articles into single Slaughter and May report")
    
    choice = input("Enter your choice (1 or 2): ").strip()
    
    if choice == "2":
        generate_single_company_report("outputs/module3_categorized_output.json", output_dir="outputs/ppt")
    else:
        generate_reports("outputs/module3_categorized_output.json", output_dir="outputs/ppt")